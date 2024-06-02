from cgi import test
from PIL import Image
import csv
import os
from os import path
import numpy as np
from itertools import count
from tkinter import filedialog
from turtle import Shape, color, width
from numpy import dtype
import matplotlib.pyplot as plt
import pandas as pd
from statistics import mean
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.dml.color import RGBColor
import time
import matplotlib.patches as patches
from matplotlib import rc
from selenium import webdriver
from selenium.webdriver.common.by import By
from urllib.request import urlopen
from bs4 import BeautifulSoup
import re
import comtypes.client
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import warnings

warnings.filterwarnings("ignore")

#Game Type (Game, Bullpen, Scrimmage)
#For Opponents leave on Scrimmage
game_type = 'Scrimmage'

#Only in use when pulling Box Score Data
home_or_away = 'Home'

#'on' or 'off' (Leave off for Opponents)
annotate = 'on'

#Folder/File name
date = "5-31-24"

#Change to Team Name in Trackman File
pitcher_team = "ORE_BEA"

split_fastballs = 'off'

in_zone = 'off'

# Opens Window to Choose CSV files
csv_files = filedialog.askopenfilenames(title="Select CSV files", filetypes=(("CSV files", "*.csv"), ("all files", "*.*")))

# Check if files were selected
if csv_files:
    # Initialize an empty list to store DataFrames
    dfs = []

    # Iterate through selected CSV files and read them into DataFrames
    for file in csv_files:
        df = pd.read_csv(file)
        dfs.append(df)

    # Concatenate all DataFrames
    global_df = pd.concat(dfs, ignore_index=True)
    global_df['TaggedPitchType'] = global_df['TaggedPitchType'].replace('FourSeamFastBall', 'Fastball')
    global_df['TaggedPitchType'] = global_df['TaggedPitchType'].replace('TwoSeamFastBall', 'Sinker')
    global_df['TaggedPitchType'] = global_df['TaggedPitchType'].replace('OneSeamFastBall', 'Sinker')
    global_df['TaggedPitchType'] = global_df['TaggedPitchType'].replace('Four-Seam', 'Fastball')
    global_df['TaggedPitchType'] = global_df['TaggedPitchType'].replace('Changeup', 'ChangeUp')
    global_df['HorzBreak'] = global_df['HorzBreak'].astype(float)
    global_df['InducedVertBreak'] = global_df['InducedVertBreak'].astype(float)




    # Print the concatenated DataFrame or perform further operations
    print(global_df)
else:
    print("No CSV files selected.")

#established global player counter
j = 0

html = 'https://osubeavers.com/sports/baseball/stats/2023/arizona-state/boxscore/20007'

# Diffrent Pitch Type catagories do not change
pitchtypes = ['Fastball', 'Fastball (R)', 'Fastball (L)','Sinker','Cutter', 'Slider','Curveball', 'ChangeUp', 'Splitter']
name_counter = 0

#List of Names for Players to Run program through
names = []
    

def name_fix_csv():
    global global_df
    global_df['Pitcher'] = global_df['Pitcher'].replace({'Lattery , AJ': 'Lattery, AJ'})
    global_df['Pitcher'] = global_df['Pitcher'].replace({'larson, Rhett': 'Larson, Rhett'})
    global_df.to_csv(csv_file, index=False)
    #global_df = global_df.drop(global_df[global_df.PitcherTeam != 'MIN_GOL'].index)
    #print(global_df)


def pull_data_from_box_score(home_or_away):
    global html
    global names
    df = pd.read_html(html)
    if home_or_away == 'Home':
        box_score_df = df[5]
    elif home_or_away == 'Away':
        box_score_df = df[4]

    box_score_df = box_score_df[:-1]
    #Removes Any Save or Win stat From BOX
    for i, row in box_score_df.iterrows():
        box_score_df.at[i,'Player'] = re.sub("[\(\[].*?[\)\]]", "", box_score_df.at[i,'Player'])
        box_score_df.at[i,'Player'] = box_score_df.at[i,'Player'].rstrip()
        names.append(box_score_df.at[i,'Player'])

    return box_score_df

def pull_names_from_trackman():
    global names
    global pitcher_team
    name_df = global_df.drop(global_df[global_df.PitcherTeam != pitcher_team].index)
    names = pd.unique(name_df['Pitcher'])
    return

def pull_stats_from_trackman():
    global global_df
    pitcher_df = global_df[global_df.Pitcher == names[name_counter]]
    pitcher_df = pitcher_df.reset_index(drop=True)
    outs = pitcher_df["OutsOnPlay"].sum() + (pitcher_df["KorBB"] == "Strikeout").sum()
    ip = str(outs//3) + '.' + str(outs%3)
    bf = 0
    for index, row in pitcher_df.iterrows(): 
        if index > 0:
            if pitcher_df.at[index, 'Batter'] != pitcher_df.at[index-1, 'Batter']:
                bf += 1
    
    hits = (pitcher_df["PlayResult"] == "Single").sum() + (pitcher_df["PlayResult"] == "Double").sum() + (pitcher_df["PlayResult"] == "Triple").sum() + (pitcher_df["PlayResult"] == "HomeRun").sum()
    walks = (pitcher_df["KorBB"] == "Walk").sum()
    hbp = (pitcher_df['PitchCall'] == 'HitByPitch').sum()
    strikeouts = (pitcher_df["KorBB"] == "Strikeout").sum()
    whip = (walks + hits)/(outs/3)
    whip = "%.3f" % round(whip,3)
    pitches = len(pitcher_df.index)
    av_ev = (pitcher_df['ExitSpeed']).mean(skipna=True)
    av_ev = "%.1f" % round(av_ev,1)

    return_list = [ip,bf,hits,walks,hbp,strikeouts,whip,av_ev,pitches]

    #print(return_list)

    return return_list


#Function pulls data from CSV in order to make the Pitch Location Charts
def get_pitch_location():

    global global_df
    pitch_loc_df = global_df
    pitch_loc_df = pitch_loc_df.drop(pitch_loc_df[pitch_loc_df.Pitcher != names[name_counter]].index)
    remove_list = pitch_loc_df.columns.values.tolist()
    remove_list.remove("Pitcher")
    remove_list.remove("PlateLocHeight")
    remove_list.remove("PlateLocSide")
    remove_list.remove("BatterSide")
    remove_list.remove('TaggedPitchType')
    #Removes all collums other than those with .remove above from data frame
    pitch_loc_df = pitch_loc_df.drop(remove_list, axis=1)
    #Switches plate loc side values to be in catcher view by multiplying by -1
    #pitch_loc_df['PlateLocSide'] = (pitch_loc_df['PlateLocSide'])
    right_loc_df = pitch_loc_df.drop(pitch_loc_df[pitch_loc_df.BatterSide != 'Right'].index)
    left_loc_df = pitch_loc_df.drop(pitch_loc_df[pitch_loc_df.BatterSide != 'Left'].index)
    return_table = []
    return_table.append(right_loc_df)
    return_table.append(left_loc_df)
    return return_table

def get_count_percents():
    global global_df
    pitcher_df = global_df.drop(global_df[global_df.Pitcher != names[name_counter]].index)
    zero_zero_df = pitcher_df[ (pitcher_df['Balls']==0) & (pitcher_df['Strikes']==0)]
    one_one_df = pitcher_df[ (pitcher_df['Balls']==1) & (pitcher_df['Strikes']==1)]
    zero_zero_pitches = len(zero_zero_df.index)
    one_one_pitches = len(one_one_df.index)
    zero_zero_strikes = (zero_zero_df['PitchCall'] == 'StrikeCalled').sum() + (zero_zero_df['PitchCall'] == 'StrikeSwinging').sum() + (zero_zero_df['PitchCall'] == 'InPlay').sum() + (zero_zero_df['PitchCall'] == 'FoulBall').sum() + (zero_zero_df['PitchCall'] == 'FoulBallNotFieldable').sum() + (one_one_df['PitchCall'] == 'FoulBallFieldable').sum()
    one_one_strikes = (one_one_df['PitchCall'] == 'StrikeCalled').sum() + (one_one_df['PitchCall'] == 'StrikeSwinging').sum() + (one_one_df['PitchCall'] == 'InPlay').sum() + (one_one_df['PitchCall'] == 'FoulBall').sum() + (one_one_df['PitchCall'] == 'FoulBallNotFieldable').sum() + (one_one_df['PitchCall'] == 'FoulBallFieldable').sum()
    zero_zero_percent = ("%.0f" % round(100*(zero_zero_strikes/zero_zero_pitches),0)) + '%'
    one_one_percent = ("%.0f" % round(100*(one_one_strikes/one_one_pitches),0)) + '%'

    return_list = [['0-0',str(zero_zero_pitches),str(zero_zero_strikes),zero_zero_percent],['1-1',str(one_one_pitches),str(one_one_strikes),one_one_percent]]

    return return_list




def get_rel_metrics():
    global global_df
    pitch_loc_df = global_df
    pitch_loc_df = pitch_loc_df.drop(pitch_loc_df[pitch_loc_df.Pitcher != names[name_counter]].index)
    remove_list = pitch_loc_df.columns.values.tolist()
    remove_list.remove("Pitcher")
    remove_list.remove("RelHeight")
    remove_list.remove("RelSide")
    remove_list.remove("Extension")
    remove_list.remove("PitcherThrows")
    remove_list.remove('TaggedPitchType')
    #Removes all collums other than those with .remove above from data frame
    pitch_loc_df = pitch_loc_df.drop(remove_list, axis=1)
    #Switches plate loc side values to be in catcher view by multiplying by -1
    #pitch_loc_df['PlateLocSide'] = (pitch_loc_df['PlateLocSide'])
    return pitch_loc_df

def new_pitch_type_tables(player_df):
    global pitchtypes
    df_list = []
    return_list = []
    for i in range(len(pitchtypes)):
        pitch_type_df = player_df[(player_df['TaggedPitchType'] == pitchtypes[i])]
        if len(pitch_type_df.index) > 0:
            df_list.append(pitch_type_df)

    for j in range(len(df_list)):
        pitch_type_list =[]
        working_df = df_list[j]
        working_df = working_df.reset_index()
        pitch_type_list.append(working_df['TaggedPitchType'].iloc[0])
        pitch_type_list.append(len(working_df.index))
        pitch_type_list.append("%.1f" % round(working_df['RelSpeed'].mean(),1))
        pitch_type_list.append("%.1f" % round(working_df['RelSpeed'].max(),1))
        pitch_type_list.append("%.0f" % round(working_df['SpinRate'].mean(),0))
        pitch_type_list.append("%.0f" % round(working_df['SpinRate'].max(),0))
        pitch_type_list.append("%.1f" % round(working_df['InducedVertBreak'].mean(),1))
        pitch_type_list.append("%.1f" % round(working_df['HorzBreak'].mean(),1))
        spin_axis_avg = working_df['SpinAxis'].mean()

        #converts spin axis to tilt
        spin_axis = spin_axis_avg
        spin_axis = float(spin_axis)
        test1 = spin_axis//30
        test2 = spin_axis%30
        test2 = test2/30
        if test1<7:
            tilt1 = test1 + 6
        elif test1>=7:
            tilt1 = test1 - 6
        tilt2 = test2*60
        tilt2 = int(round(tilt2,0))
        if tilt2<10:
            tilt2str = str('0' + str(tilt2))
        else:
            tilt2str = str(tilt2)

        pitch_type_list.append(str(int(round(tilt1,0))) + ':' + tilt2str)
        cleaned_df = working_df.dropna(subset=['PlateLocSide', 'PlateLocHeight'])
        in_zone_df = working_df[(working_df['PlateLocSide'] < 0.83083) & (working_df['PlateLocSide'] > -0.83083) & (working_df['PlateLocHeight'] < 3.67333) & (working_df['PlateLocHeight'] > 1.52417)]
        #vaa_df = working_df[(working_df['PlateLocSide'] < 0.83083) & (working_df['PlateLocSide'] > -0.83083) & (working_df['PlateLocHeight'] < 2.24056) & (working_df['PlateLocHeight'] > 1.52417)]
        pitch_type_list.append("%.2f" % round(in_zone_df['VertApprAngle'].mean(),2))
        pitch_type_list.append("%.2f" % round(working_df['RelHeight'].mean(),2))
        pitch_type_list.append("%.2f" % round(working_df['RelSide'].mean(),2))
        pitch_type_list.append("%.2f" % round(working_df['Extension'].mean(),2))

        strikes = 0
        swings = 0
        whiffs = 0
        for d in range(len(working_df)):
            if working_df.at[d,'PitchCall'] in ('FoulBall', 'StrikeCalled', 'StrikeSwinging', 'InPlay', 'FoulBallNotFieldable', 'FoulBallFieldable'):
                strikes = strikes + 1
            if working_df.at[d,'PitchCall'] in ('FoulBall', 'StrikeSwinging', 'InPlay','FoulBallFieldable', 'FoulBallNotFieldable'):
                swings = swings + 1
            if working_df.at[d,'PitchCall'] == 'StrikeSwinging':
                whiffs = whiffs + 1

        if in_zone == 'off':
            strike_percentage = round(100*(strikes/len(working_df)),0)
            pitch_type_list.append(str(str(int(strike_percentage)) + '%'))
        else:
            in_zone_percentage = round(100*(len(in_zone_df)/len(cleaned_df)),0)
            pitch_type_list.append(str(str(int(in_zone_percentage)) + '%'))

        if swings > 0:
            whiff_percentage = round(100*(whiffs/swings),0)
            pitch_type_list.append(str(str(int(whiff_percentage)) + '%'))
        else:
            pitch_type_list.append('N/A')
        
        string_list = [str(value) for value in pitch_type_list]

        return_list.append(string_list)

    return return_list
        



#Create Chart Based on Pitch Location
def pitch_loaction_chart(player_df):
    global annotate
    plt.rc('axes',edgecolor='w')
    sides = ["Right","Left"]
    player_df['counter'] = range(len(player_df))
    table = [player_df[player_df['BatterSide'] == 'Right'],player_df[player_df['BatterSide'] == 'Left']]
    table[0].reset_index(inplace=True)
    table[1].reset_index(inplace=True)

    #Loop runs once for LHH locations and once for RHH locations
    for j in range(2):
        if j == 0:
            img = plt.imread("rhh.png")
        if j == 1:
            img = plt.imread("lhh.png")
        fig, ax = plt.subplots()
        ax.imshow(img, extent=[-2.57,2.57,-0.35,5.30])
        #rect = patches.Rectangle((-0.708333, 1.6466667), 1.4166667, 1.90416667, linewidth=1, edgecolor='black', facecolor='none')
        #ax.add_patch(rect)
        #Plots each pitch type a diffrent color
        for i in range(len(table[j].index)):
            if table[j].loc[i, 'PlateLocSide'] or table[j].loc[i, 'PlateLocHeight'] != None:
                if table[j].loc[i, 'TaggedPitchType'] == 'Fastball':
                    ax.plot(table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight'], color="#eb0e0e", marker='o', linestyle='', markersize=8)
                    if annotate == 'on':
                        plt.annotate(str(table[j].loc[i, 'counter']+1), (table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight']), textcoords="offset points", xytext=(0,6), ha='center')
                elif table[j].loc[i, 'TaggedPitchType'] == 'Sinker':
                    ax.plot(table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight'], color="#fc535f", marker='o', linestyle='', markersize=8)
                    if annotate == 'on':
                        plt.annotate(str(table[j].loc[i, 'counter']+1), (table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight']), textcoords="offset points", xytext=(0,6), ha='center')
                elif table[j].loc[i, 'TaggedPitchType'] == 'Cutter':
                    ax.plot(table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight'], color="#fcbe03", marker='o', linestyle='', markersize=8)
                    if annotate == 'on':
                        plt.annotate(str(table[j].loc[i, 'counter']+1), (table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight']), textcoords="offset points", xytext=(0,6), ha='center')
                elif table[j].loc[i, 'TaggedPitchType'] == 'Slider':
                    ax.plot(table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight'], color="#1105f7", marker='o', linestyle='', markersize=8)
                    if annotate == 'on':
                        plt.annotate(str(table[j].loc[i, 'counter']+1), (table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight']), textcoords="offset points", xytext=(0,6), ha='center')
                elif table[j].loc[i, 'TaggedPitchType'] == 'Curveball':
                    ax.plot(table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight'], color="#05bff7", marker='o', linestyle='', markersize=8)
                    if annotate == 'on':
                        plt.annotate(str(table[j].loc[i, 'counter']+1), (table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight']), textcoords="offset points", xytext=(0,6), ha='center')
                elif table[j].loc[i, 'TaggedPitchType'] in ('ChangeUp', 'Splitter'):
                    ax.plot(table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight'], color="#f763f0", marker='o', linestyle='', markersize=8)
                    if annotate == 'on':
                        plt.annotate(str(table[j].loc[i, 'counter']+1), (table[j].loc[i, 'PlateLocSide'], table[j].loc[i, 'PlateLocHeight']), textcoords="offset points", xytext=(0,6), ha='center')
        #Creates path and saves it as photo
        ax.tick_params(axis = 'x', colors = 'white')
        ax.tick_params(axis = 'y', colors = 'white')
        plt.xlim([-2.57, 2.57])
        plt.ylim([-0.5, 5.4])
        plt.savefig(os.path.join(date, names[name_counter], sides[j] + '.png'),bbox_inches='tight', pad_inches = 0)

#Creates Breakplot
def breakplot(player_df):
    player_df['counter'] = range(len(player_df))
    player_df.reset_index(inplace=True)
    fig, bp = plt.subplots()
    # Move left y-axis and bottim x-axis to centre, passing through (0,0)
    bp.spines['left'].set_position('center')
    bp.spines['bottom'].set_position('center')
    bp.spines['left'].set_color('#8c8c8c')
    bp.spines['bottom'].set_color('#8c8c8c')

    # Eliminate upper and right axes
    bp.spines['right'].set_color('none')
    bp.spines['top'].set_color('none')

    bp.tick_params(axis='x', colors='#8c8c8c')   # Set the color of x-axis ticks
    bp.tick_params(axis='y', colors='#8c8c8c') 

    # Show ticks in the left and lower axes only
    bp.xaxis.set_ticks_position('bottom')
    bp.yaxis.set_ticks_position('left')
    for i in range(len(player_df)):
        pitch_type = player_df.loc[i, 'TaggedPitchType']
        if pitch_type == 'Fastball':
            bp.plot(player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak'], color="#eb0e0e", marker='o', linestyle='', markersize=6)
            if annotate == 'on':
                plt.annotate(str(player_df.loc[i, 'counter']+1), (player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak']), textcoords="offset points", xytext=(0,4), ha='center', fontsize=6, zorder=10)
        if pitch_type == 'Sinker':
            bp.plot(player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak'], color="#fc535f", marker='o', linestyle='', markersize=6)
            if annotate == 'on':
                plt.annotate(str(player_df.loc[i, 'counter']+1), (player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak']), textcoords="offset points", xytext=(0,4), ha='center', fontsize=6, zorder=10)
        elif pitch_type == 'Cutter':
            bp.plot(player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak'], color="#fcbe03", marker='o', linestyle='', markersize=6)
            if annotate == 'on':
                plt.annotate(str(player_df.loc[i, 'counter']+1), (player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak']), textcoords="offset points", xytext=(0,4), ha='center', fontsize=6, zorder=10)
        elif pitch_type == 'Slider':
            bp.plot(player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak'], color="#1105f7", marker='o', linestyle='', markersize=6)
            if annotate == 'on':
                plt.annotate(str(player_df.loc[i, 'counter']+1), (player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak']), textcoords="offset points", xytext=(0,4), ha='center', fontsize=6, zorder=10)
        elif pitch_type == 'Curveball':
            bp.plot(player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak'], color="#05bff7", marker='o', linestyle='', markersize=6)
            if annotate == 'on':
                plt.annotate(str(player_df.loc[i, 'counter']+1), (player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak']), textcoords="offset points", xytext=(0,4), ha='center', fontsize=6, zorder=10)
        elif pitch_type in ('ChangeUp', 'Splitter'):
            bp.plot(player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak'], color="#f763f0", marker='o', linestyle='', markersize=6)
            if annotate == 'on':
                plt.annotate(str(player_df.loc[i, 'counter']+1), (player_df.loc[i, 'HorzBreak'], player_df.loc[i, 'InducedVertBreak']), textcoords="offset points", xytext=(0,4), ha='center', fontsize=6, zorder=10)
    plt.xlim([-27, 27])
    plt.ylim([-27, 27])


    newpath = os.path.join(date, names[name_counter])
    if not os.path.exists(newpath):
        os.makedirs(newpath)

    plt.savefig(os.path.join(date, names[name_counter], 'breakplot' + '.png'))              
      
    return
#Creates Batter View Release Point Chart
def rel_batter_view(pitch_type_tables):
    fig, ax = plt.subplots()
    side = pitch_type_tables.iat[0,1]
    pitch_type_tables['RelSide'] = (pitch_type_tables['RelSide'] * -1)
    if side == 'Right':
        img = plt.imread("Rel3.png")
        ax.imshow(img, extent=[-4.5,2.2,0,8])
    elif side == 'Left':
        img = plt.imread('Rel2.png')
        ax.imshow(img, extent=[-2.2,4.5,0,8])
    
    for j in range(len(pitch_type_tables.index)):
        pitch_type = pitch_type_tables.iat[j,2]
        if pitch_type == 'Fastball':
            ax.plot(pitch_type_tables.iat[j,4], pitch_type_tables.iat[j,3], color="#eb0e0e", marker='o', linestyle='', markersize=5)
        if pitch_type == 'Sinker':
            ax.plot(pitch_type_tables.iat[j,4], pitch_type_tables.iat[j,3], color="#fc535f", marker='o', linestyle='', markersize=5)
        elif pitch_type == 'Cutter':
            ax.plot(pitch_type_tables.iat[j,4], pitch_type_tables.iat[j,3], color="#fcbe03", marker='o', linestyle='', markersize=5)
        elif pitch_type == 'Slider':
            ax.plot(pitch_type_tables.iat[j,4], pitch_type_tables.iat[j,3], color="#1105f7", marker='o', linestyle='', markersize=5)
        elif pitch_type == 'Curveball':
            ax.plot(pitch_type_tables.iat[j,4], pitch_type_tables.iat[j,3], color="#05bff7", marker='o', linestyle='', markersize=5)
        elif pitch_type in ('ChangeUp', 'Splitter'):
            ax.plot(pitch_type_tables.iat[j,4], pitch_type_tables.iat[j,3], color="#f763f0", marker='o', linestyle='', markersize=5)


    newpath = os.path.join(date, names[name_counter])
    if not os.path.exists(newpath):
        os.makedirs(newpath)
    plt.tick_params(left = False, bottom = False)
    ax.axes.xaxis.set_ticklabels([])
    ax.axes.yaxis.set_ticklabels([])
    plt.savefig(os.path.join(date, names[name_counter], 'RelH_BV' + '.png'), bbox_inches='tight', pad_inches = -0.02)              
      
    return
#
def rel_90_view(pitch_type_tables):
    fig, ax = plt.subplots()
    img = plt.imread("Rel90_2.jpg")
    ax.imshow(img, extent=[-7,10,0,8])
    
    for j in range(len(pitch_type_tables.index)):
        pitch_type = pitch_type_tables.iat[j,2]
        if pitch_type == 'Fastball':
            ax.plot(pitch_type_tables.iat[j,5], pitch_type_tables.iat[j,3], color="#eb0e0e", marker='o', linestyle='', markersize=5)
        if pitch_type == 'Sinker':
            ax.plot(pitch_type_tables.iat[j,5], pitch_type_tables.iat[j,3], color="#fc535f", marker='o', linestyle='', markersize=5)
        elif pitch_type == 'Cutter':
            ax.plot(pitch_type_tables.iat[j,5], pitch_type_tables.iat[j,3], color="#fcbe03", marker='o', linestyle='', markersize=5)
        elif pitch_type == 'Slider':
            ax.plot(pitch_type_tables.iat[j,5], pitch_type_tables.iat[j,3], color="#1105f7", marker='o', linestyle='', markersize=5)
        elif pitch_type == 'Curveball':
            ax.plot(pitch_type_tables.iat[j,5], pitch_type_tables.iat[j,3], color="#05bff7", marker='o', linestyle='', markersize=5)
        elif pitch_type in ('ChangeUp', 'Splitter'):
            ax.plot(pitch_type_tables.iat[j,5], pitch_type_tables.iat[j,3], color="#f763f0", marker='o', linestyle='', markersize=5)

    newpath = os.path.join(date, names[name_counter])
    if not os.path.exists(newpath):
        os.makedirs(newpath)
    plt.tick_params(left = False, bottom = False)
    ax.axes.xaxis.set_ticklabels([])
    ax.axes.yaxis.set_ticklabels([])
    plt.savefig(os.path.join(date, names[name_counter], 'RelH_90V' + '.png'), bbox_inches='tight', pad_inches = -0.02)              
    return

'''def spin_efficency(player_df):
    player_df['yR'] = 60.5 - player_df['Extension']
    player_df['tR'] = (-player_df['vy0']-(player_df['vy0']**2-2*player_df['ay0']*(50-player_df['yR']))**0.5)/player_df['ay0']
    player_df['vxR'] = 
'''


#Circular Bar Plot Data Formatting
def get_bar_plot_data():
    global global_df
    df = global_df.drop((global_df[global_df.Pitcher != names[name_counter]].index))
    remove_list = df.columns.values.tolist()
    remove_list.remove("SpinAxis")
    remove_list.remove("TaggedPitchType")
    #Removes all collums other than those with .remove above from data frame
    df = df.drop(remove_list, axis=1)
    #for i in range(len(df)):
    #    df.



#Stacked Circular BarPlot
def stacked_bar_plot(df):
    ANGLES = np.linspace(0, 2 * np.pi, len(df), endpoint=False)
    VALUES = df["value"].values
    LABELS = df["name"].values

    # Determine the width of each bar. 
    # The circumference is '2 * pi', so we divide that total width over the number of bars.
    WIDTH = 2 * np.pi / len(VALUES)

    # Determines where to place the first bar. 
    # By default, matplotlib starts at 0 (the first bar is horizontal)
    # but here we say we want to start at pi/2 (90 deg)
    OFFSET = np.pi / 2

    # Initialize Figure and Axis
    fig, ax = plt.subplots(figsize=(20, 10), subplot_kw={"projection": "polar"})

    # Specify offset
    ax.set_theta_offset(OFFSET)

    # Set limits for radial (y) axis. The negative lower bound creates the whole in the middle.
    ax.set_ylim(-100, 100)

    # Remove all spines
    ax.set_frame_on(False)

    # Remove grid and tick marks
    ax.xaxis.grid(False)
    ax.yaxis.grid(False)
    ax.set_xticks([])
    ax.set_yticks([])

    # Add bars
    ax.bar(
        ANGLES, VALUES, width=WIDTH, linewidth=2,
        color="#61a4b2", edgecolor="white"
    )

#Converts the Power Point to a PDF
def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()

def stats_from_box_score(box_score_df):
    list = box_score_df.loc[name_counter, :].values.tolist()
    list.pop(0) #Removes Name
    list.pop(7) #Removes BK
    list.pop(8) #Removes IBB
    list.pop(8) #Removes AB
    list.pop(9) #Removes FO
    list.pop(9) #Removes GO
    list.insert(1, list.pop(8)) #Moves BF to pos 1
    list.insert(4, list.pop(2)) #Moves hits to pos 4
    list.insert(6, list.pop(8)) #Moves HBP to pos 6
    list.insert(7, list.pop(8)) #Swaps WP and SO
    whip = (list[4]+list[5])/((list[0]//1)+(list[0]%1)*3.3333333)
    whip = "%.3f" % round(whip,3)
    list.insert(9, whip)
    return list

def get_bip_info():
    global global_df
    bip_df = global_df.drop((global_df[global_df.Pitcher != names[name_counter]].index))
    bip_df = bip_df.drop((bip_df[bip_df.PitchCall != 'InPlay'].index))
    bips = str(len(bip_df.index))
    bip_df = bip_df.dropna(subset = ['ExitSpeed'])
    count = bip_df.count()
    if count['ExitSpeed'] > 0:
        avg_ev = str(round(mean(bip_df['ExitSpeed']),1))
    else:
        avg_ev = 'N/A'
    hard_hits = (bip_df['ExitSpeed'] > 90.0).sum()
    measuerd_bips_ev = count['ExitSpeed']
    measuerd_bips_la = count["Angle"]
    hard_hit_percentage = ("%.1f" % round(100*hard_hits/measuerd_bips_ev,1)) +'%'
    grounders = (bip_df['Angle'] < 10.0).sum()
    grounder_percentage = ("%.1f" % round(100*grounders/measuerd_bips_la,1)) + '%'
    return_list = [bips,avg_ev,hard_hit_percentage,grounder_percentage]
    return return_list



# Does the Presentation work through PPTX to create the presentation
def create_presentation(averages):
    arsenal = len(averages)
    # ---create presentation with 1 slide---
    prs = Presentation("Template2.pptx")
    prs.slide_width = Inches(11)
    prs.slide_height = Inches(8.5)
    #slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide = prs.slides.get(257)

    name = names[name_counter].split()
    #print(name)
    full_name = name[-1] + ' ' + name[0]
    full_name = full_name[:-1]
    title = slide.shapes.title
    title.text = full_name
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.name = 'Beaver Bold'
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ---add table to slide---
    x, y, cx, cy = Inches(0.15), Inches(1), Inches(10.6), Inches(2.5)
    shape = slide.shapes.add_table((arsenal + 1), 12, x, y, cx, cy)
    table = shape.table
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(0.65)

    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
    tbl[0][-1].text = style_id

    table_labels = ['Pitch Type','#',"AVG Velo","MAX Velo","AVG Spin","MAX Spin","Vert Break","Horz Break",'Tilt',"Rel Height","Rel Side","Exten-sion","InZone%","Whiff%"]

    #creating labels for values in all tables
    for i in range(len(table_labels)-2):
        cell = table.cell(0, i)
        cell.text = table_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        if (i == 9 or i==12):
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    #nested for loop which puts vales from averages into table
    for i in range(arsenal):
        for j in range(12):
            cell = table.cell((i+1),j)
            cell.text = averages[i][j]
            cell.text_frame.paragraphs[0].font.size = Pt(15)
            cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
    prs.save(os.path.join(date, names[name_counter], names[name_counter] + '.pptx'))

    imbp = Image.open(os.path.join(date, names[name_counter], 'breakplot' + '.png'))
    imbp.crop((65,50,590,435)).save(os.path.join(date, names[name_counter], 'breakplot' + '.png'))

    img_bp = os.path.join(date, names[name_counter], 'breakplot' + '.png')
    add_picture = slide.shapes.add_picture(img_bp,Inches(6.3866666),Inches(3.8), width=Inches(4.4), height=Inches(3.23))
    line = add_picture.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Inches(0.025)

        # Opens a image in RGB mode
    im = Image.open(os.path.join(date, names[name_counter], 'Right' + '.png'))

    # Size of the image in pixels (size of original image)
    # (This is not mandatory)
    widthr, heightr = im.size
 
    # Cropped image of above dimension
    # (It will not change original image)
    im1 = im.crop((20,6,widthr,369)).save(os.path.join(date, names[name_counter], 'Right' + '.png'))

    img_rhh = os.path.join(date, names[name_counter], 'Right' + '.png')
    add_picture2 = slide.shapes.add_picture(img_rhh,Inches(3.3),Inches(3.8), width=Inches(2.8474), height=Inches(3.23))
    line = add_picture2.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Inches(0.025)

    iml = Image.open(os.path.join(date, names[name_counter], 'Left' + '.png'))
    widthl, heightl = iml.size
    im2 = iml.crop((20,6,widthl,369)).save(os.path.join(date, names[name_counter], 'Left' + '.png'))

    img_lhh = os.path.join(date, names[name_counter], 'Left' + '.png')
    add_picture3 = slide.shapes.add_picture(img_lhh,Inches(.20),Inches(3.8), width=Inches(2.8474), height=Inches(3.23))
    line = add_picture3.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Inches(0.025)

    newpath = os.path.join(date, 'PDFs')
    if not os.path.exists(newpath):
        os.makedirs(newpath)

    prs.save(os.path.join(date, names[name_counter], names[name_counter] + '.pptx'))

    PPTtoPDF(os.path.join(date, names[name_counter], names[name_counter] + '.pptx'),os.path.join(date, 'PDFs', name[0][:-1] + ' ' + date + '.pdf'))

    return


# Does the Presentation work through PPTX to create the presentation
def create_presentation_game(averages,bip_info,count_info):
    #print(stats)
    arsenal = len(averages)
    # ---create presentation with 1 slide---
    prs = Presentation("Template2.pptx")
    prs.slide_width = Inches(11)
    prs.slide_height = Inches(8.5)
    #slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide = prs.slides.get(257)

    name = names[name_counter].split()
    print(name)
    full_name = name[-1] + ' ' + name[0]
    full_name = full_name[:-1]
    title = slide.shapes.title
    title.text = full_name
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.name = 'Beaver Bold'
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ---add table to slide---
    if arsenal > 4:
        y_value = 1
    else:
        y_value = 1.16
    x, y, cx, cy = Inches(0.15), Inches(y_value), Inches(10.3), Inches(2.6)
    #y is 1.16 when switching back to game
    shape = slide.shapes.add_table((arsenal + 1), 15, x, y, cx, cy)
    table = shape.table
    table.columns[0].width = Inches(1.15)
    table.columns[1].width = Inches(0.58)

    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
    tbl[0][-1].text = style_id

    if in_zone == 'off':
        table_labels = ['Pitch Type','#',"AVG Velo","MAX Velo","AVG Spin","MAX Spin","Vert Break",
        "Horz Break",'Tilt','InZone VAA',"Rel Height","Rel Side","Exten-sion","Strike%","Whiff%"]
    else: 
        table_labels = ['Pitch Type','#',"AVG Velo","MAX Velo","AVG Spin","MAX Spin","Vert Break",
        "Horz Break",'Tilt','InZone VAA',"Rel Height","Rel Side","Exten-sion","InZone%","Whiff%"]


    #creating labels for values in all tables
    for i in range(len(table_labels)):
        cell = table.cell(0, i)
        cell.text = table_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(13)
        if (i == 10 or i==13):
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        if i==9:
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        if i==13:
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    #nested for loop which puts vales from averages into table
    for i in range(arsenal):
        for j in range(15):
            cell = table.cell((i+1),j)
            cell.text = averages[i][j]
            cell.text_frame.paragraphs[0].font.size = Pt(14.5)
            if j == 9:
                cell.text_frame.paragraphs[0].font.size = Pt(12)
            if averages[i][j] == 'Fastball (R)' or averages[i][j] == 'Fastball (L)':
                cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
    prs.save(os.path.join(date, names[name_counter], names[name_counter] + '.pptx'))

    if game_type == 'Game':
        top_table_labels = ['IP','BF','R','ER','H','BB','HBP','WP','SO','WHIP','Pitches']
    elif game_type == 'Scrimmage':
        top_table_labels = ['IP','BF','H','BB','HBP','SO','WHIP','AVG EV','Pitches']
    else:
        top_table_lables = []

    if game_type == 'Game':
        x, y, cx, cy = Inches(0.15), Inches(0.45), Inches(10.675), Inches(0.63)
        shape = slide.shapes.add_table(2, len(top_table_labels), x, y, cx, cy)
        table = shape.table

        tbl =  shape._element.graphic.graphicData.tbl
        style_id = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
        tbl[0][-1].text = style_id

        for i in range(len(top_table_labels)):
            cell = table.cell(0, i)
            cell.text = top_table_labels[i]
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        '''
        for j in range(len(stats)):
            cell = table.cell(1,j)
            stats[j] = str(stats[j])
            cell.text = stats[j]
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER 
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        '''

    hit_table_labels = ["Balls in Play", "AVG EV", 'Hard Hit %', 'Grounder %']
    if game_type != 'Bullpen':

        x, y, cx, cy = Inches(3.4), Inches(6.54), Inches(2.1), Inches(1.81)    
        shape = slide.shapes.add_table(4, 2, x, y, cx, cy)
        table = shape.table

        tbl =  shape._element.graphic.graphicData.tbl
        style_id = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
        tbl[0][-1].text = style_id
        table.columns[0].width = Inches(1.3)

        for i in range(4):
            cell = table.cell(i, 0)
            cell.text = hit_table_labels[i]
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(231, 231, 231)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for i in range(4):
            cell = table.cell(i,1)
            cell.text = bip_info[i]
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(231, 231, 231)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        x, y, cx, cy = Inches(5.83), Inches(6.54), Inches(2.88), Inches(1.81)    
        shape = slide.shapes.add_table(3, 4, x, y, cx, cy)
        table = shape.table

        tbl =  shape._element.graphic.graphicData.tbl
        style_id = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
        tbl[0][-1].text = style_id

        count_labels = ['Count','Pitches','Strikes','Strike%']

        for i in range(4):
            cell = table.cell(0, i)
            cell.text = count_labels[i]
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        for j in range(2):
            for i in range(4):
                cell = table.cell(j+1, i)
                cell.text = count_info[j][i]
                cell.text_frame.paragraphs[0].font.size = Pt(14)
                cell.text_frame.paragraphs[0].font.name = 'Bahnschrift'
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE



    imbp = Image.open('RelH.jpg')
    widthr, heightr = imbp.size
    imbp.crop((195,233,widthr,heightr)).save('RelH2.jpg')

    imbp = Image.open(os.path.join(date, names[name_counter], 'breakplot' + '.png'))
    imbp.crop((65,50,590,435)).save(os.path.join(date, names[name_counter], 'breakplot' + '.png'))

    img_bp = os.path.join(date, names[name_counter], 'breakplot' + '.png')
    add_picture = slide.shapes.add_picture(img_bp,Inches(7.24),Inches(3.83), width=Inches(3.58), height=Inches(2.63))
    line = add_picture.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Inches(0.015)


    # Opens a image in RGB mode
    im = Image.open(os.path.join(date, names[name_counter], 'Right' + '.png'))

    # Size of the image in pixels (size of original image)
    # (This is not mandatory)
    widthr, heightr = im.size
 
    # Cropped image of above dimension
    # (It will not change original image)
    im1 = im.crop((20,6,widthr,369)).save(os.path.join(date, names[name_counter], 'Right' + '.png'))

    img_rhh = os.path.join(date, names[name_counter], 'Right' + '.png')
    add_picture2 = slide.shapes.add_picture(img_rhh,Inches(4.85),Inches(3.83), width=Inches(2.32), height=Inches(2.63))
    line = add_picture2.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Inches(0.015)

    iml = Image.open(os.path.join(date, names[name_counter], 'Left' + '.png'))
    widthl, heightl = iml.size
    im2 = iml.crop((20,6,widthl,369)).save(os.path.join(date, names[name_counter], 'Left' + '.png'))

    img_lhh = os.path.join(date, names[name_counter], 'Left' + '.png')
    add_picture3 = slide.shapes.add_picture(img_lhh,Inches(2.46),Inches(3.83), width=Inches(2.32), height=Inches(2.63))
    line = add_picture3.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Inches(0.015)

    img_hei = os.path.join(date, names[name_counter], 'RelH_BV' + '.png')
    add_picture4 = slide.shapes.add_picture(img_hei,Inches(0.15),Inches(3.83), width=Inches(2.23), height=Inches(2.63))
    line = add_picture4.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Inches(0.015)

    img_hei90 = os.path.join(date, names[name_counter], 'RelH_90V' + '.png')
    add_picture5 = slide.shapes.add_picture(img_hei90,Inches(0.15),Inches(6.53), width=Inches(3.17), height=Inches(1.81))
    line = add_picture5.line
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Inches(0.015)

    newpath = os.path.join(date, 'PDFs')
    if not os.path.exists(newpath):
        os.makedirs(newpath)

    prs.save(os.path.join(os.getcwd(), date, names[name_counter], names[name_counter] + '.pptx'))

    PPTtoPDF(os.path.join(os.getcwd(), date, names[name_counter], names[name_counter] + '.pptx'),os.path.join(os.getcwd(), date, 'PDFs', name[0][:-1] + name[1] + ' ' + date + '.pdf'))

    return
#Main

def split_fastballs_func(player_df):
    player_df.loc[(player_df['TaggedPitchType'] == 'Fastball') & (player_df['BatterSide'] == 'Left'), 'TaggedPitchType'] = 'Fastball (L)'
    player_df.loc[(player_df['TaggedPitchType'] == 'Fastball') & (player_df['BatterSide'] == 'Right'), 'TaggedPitchType'] = 'Fastball (R)'
    return player_df

def main ():
    global names
    global home_or_away
    global name_counter
    global split_fastballs
    #name_fix_csv()
    pull_names_from_trackman()
    #box = pull_data_from_box_score(home_or_away)
    for k in range(len(names)):
        print(names[name_counter])
        player_df = global_df.drop(global_df[global_df.Pitcher != names[name_counter]].index)
        #pitch_tables = get_pitch_type_tables()
        #df = get_pitch_location()
        rel = get_rel_metrics()
        breakplot(player_df)
        rel_batter_view(rel)
        rel_90_view(rel)
        pitch_loaction_chart(player_df)
        #create_presentation_game(average_pitch_types(),stats_from_box_score(box),get_bip_info(), get_count_percents())
        #create_presentation_game(average_pitch_types(),pull_stats_from_trackman(),get_bip_info(), get_count_percents())
        #create_presentation_game(average_pitch_types(),get_bip_info(), get_count_percents())
        if split_fastballs == 'on':
            player_df = split_fastballs_func(player_df)
        create_presentation_game(new_pitch_type_tables(player_df),get_bip_info(), get_count_percents())
        name_counter = name_counter + 1



main()
